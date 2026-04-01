"""
Kinto Global — Report Artifact Generator v2
Renders branded DOCX and PPTX from rich consulting-narrative JSON payloads.
Produces professional advisory reports, not dashboard summaries.
"""
import argparse, json, textwrap
from pathlib import Path
from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches as Pi, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colours ─────────────────────────────────────────────────────────────
TEAL     = RGBColor(26, 188, 176)
DARK     = RGBColor(15, 25, 35)
MUTED    = RGBColor(100, 116, 139)
DANGER   = RGBColor(220, 38, 38)
WARN     = RGBColor(202, 138, 4)
SUCCESS  = RGBColor(5, 150, 105)
WHITE    = RGBColor(255, 255, 255)
LIGHT    = RGBColor(248, 250, 252)
MID_GREY = RGBColor(226, 232, 240)

P_TEAL   = PRGBColor(26, 188, 176)
P_DARK   = PRGBColor(15, 25, 35)
P_MUTED  = PRGBColor(100, 116, 139)
P_DANGER = PRGBColor(220, 38, 38)
P_WARN   = PRGBColor(202, 138, 4)
P_WHITE  = PRGBColor(255, 255, 255)


def s(v, fallback='—'):
    if v is None: return fallback
    t = str(v).strip()
    return t if t else fallback


def ensure_dir(p):
    Path(p).parent.mkdir(parents=True, exist_ok=True)


def priority_color(val):
    v = str(val or '').lower()
    if v in ('critical', 'high'):   return DANGER
    if v in ('medium', 'moderate'): return WARN
    return SUCCESS


def p_priority_color(val):
    v = str(val or '').lower()
    if v in ('critical', 'high'):   return P_DANGER
    if v in ('medium', 'moderate'): return P_WARN
    return P_TEAL


def score_color(pct):
    if pct < 40: return DANGER
    if pct < 65: return WARN
    return TEAL


def p_score_color(pct):
    if pct < 40: return P_DANGER
    if pct < 65: return P_WARN
    return P_TEAL


# ── DOCX helpers ──────────────────────────────────────────────────────────────
def set_cell_bg(cell, hex_str):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), hex_str.replace('#', ''))
    tcPr.append(shd)


def add_divider(doc, color='E2E8F0'):
    p = doc.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    pBdr = OxmlElement('w:pBdr')
    bottom = OxmlElement('w:bottom')
    bottom.set(qn('w:val'), 'single')
    bottom.set(qn('w:sz'), '4')
    bottom.set(qn('w:space'), '1')
    bottom.set(qn('w:color'), color)
    pBdr.append(bottom)
    pPr.append(pBdr)
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after = Pt(8)


def section_tag(doc, tag):
    p = doc.add_paragraph()
    r = p.add_run(tag.upper())
    r.bold = True
    r.font.size = Pt(8)
    r.font.color.rgb = TEAL
    r.font.name = 'Calibri'
    p.paragraph_format.space_before = Pt(18)
    p.paragraph_format.space_after = Pt(2)


def h1(doc, text, color=None):
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.bold = True
    r.font.size = Pt(18)
    r.font.color.rgb = color or DARK
    r.font.name = 'Calibri'
    p.paragraph_format.space_before = Pt(4)
    p.paragraph_format.space_after = Pt(6)
    return p


def h2(doc, text, color=None):
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.bold = True
    r.font.size = Pt(13)
    r.font.color.rgb = color or DARK
    r.font.name = 'Calibri'
    p.paragraph_format.space_before = Pt(12)
    p.paragraph_format.space_after = Pt(4)
    return p


def h3(doc, text, color=None):
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.bold = True
    r.font.size = Pt(11)
    r.font.color.rgb = color or DARK
    r.font.name = 'Calibri'
    p.paragraph_format.space_before = Pt(8)
    p.paragraph_format.space_after = Pt(3)
    return p


def body(doc, text, color=None, size=10.5):
    if not text or text == '—': return
    p = doc.add_paragraph()
    r = p.add_run(str(text))
    r.font.size = Pt(size)
    r.font.color.rgb = color or MUTED
    r.font.name = 'Calibri'
    p.paragraph_format.space_after = Pt(5)
    p.paragraph_format.line_spacing = Pt(14)
    return p


def bullet(doc, text, color=None, indent=False):
    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Inches(0.3 if indent else 0.2)
    r = p.add_run(f'• {text}')
    r.font.size = Pt(10)
    r.font.color.rgb = color or MUTED
    r.font.name = 'Calibri'
    p.paragraph_format.space_after = Pt(3)


def labelled(doc, label, value, label_color=None, value_color=None):
    p = doc.add_paragraph()
    k = p.add_run(f'{label}:  ')
    k.bold = True; k.font.size = Pt(10); k.font.color.rgb = label_color or DARK; k.font.name = 'Calibri'
    v = p.add_run(str(value))
    v.font.size = Pt(10); v.font.color.rgb = value_color or MUTED; v.font.name = 'Calibri'
    p.paragraph_format.space_after = Pt(3)


def shaded_block(doc, label, content, bg='F0FDFA', border='1ABCB0'):
    """A shaded content block — used for findings/recommendations/actions."""
    tbl = doc.add_table(rows=1, cols=1)
    tbl.style = 'Table Grid'
    cell = tbl.rows[0].cells[0]
    set_cell_bg(cell, bg)
    # Label
    lp = cell.paragraphs[0]
    lr = lp.add_run(label)
    lr.bold = True; lr.font.size = Pt(9.5); lr.font.color.rgb = TEAL; lr.font.name = 'Calibri'
    lp.paragraph_format.space_after = Pt(2)
    # Content
    cp = cell.add_paragraph()
    cr = cp.add_run(str(content))
    cr.font.size = Pt(10); cr.font.color.rgb = DARK; cr.font.name = 'Calibri'
    cp.paragraph_format.space_after = Pt(2)
    doc.add_paragraph().paragraph_format.space_after = Pt(4)


def score_table(doc, scores):
    """Module/domain score table."""
    tbl = doc.add_table(rows=1, cols=5)
    tbl.style = 'Table Grid'
    for cell, lbl in zip(tbl.rows[0].cells, ['Module / Domain', 'Score', 'Maturity Band', 'Priority', 'Status']):
        set_cell_bg(cell, '0F1923')
        r = cell.paragraphs[0].add_run(lbl)
        r.bold = True; r.font.size = Pt(8.5); r.font.color.rgb = WHITE; r.font.name = 'Calibri'
    for row_data in scores:
        row = tbl.add_row()
        pct_val = float(row_data.get('score') or row_data.get('percentage') or 0)
        name = s(row_data.get('name') or row_data.get('module_name') or row_data.get('area_name'))
        bnd  = s(row_data.get('score_band') or row_data.get('band'))
        pri  = s(row_data.get('priority') or ('Critical' if pct_val < 40 else 'High' if pct_val < 65 else 'Medium'))
        status = 'Critical' if pct_val < 40 else 'Developing' if pct_val < 65 else 'Strong'
        cells = row.cells
        cells[0].paragraphs[0].add_run(name).font.size = Pt(9)
        pr = cells[1].paragraphs[0].add_run(f'{pct_val:.1f}%')
        pr.bold = True; pr.font.size = Pt(9); pr.font.color.rgb = score_color(pct_val)
        cells[2].paragraphs[0].add_run(bnd).font.size = Pt(9)
        prir = cells[3].paragraphs[0].add_run(pri)
        prir.font.size = Pt(9); prir.font.color.rgb = priority_color(pri)
        sr = cells[4].paragraphs[0].add_run(status)
        sr.font.size = Pt(9); sr.font.color.rgb = score_color(pct_val)
    doc.add_paragraph().paragraph_format.space_after = Pt(6)


def finding_block(doc, finding, number=None):
    num_str = f"Finding {finding.get('number') or number}: " if (finding.get('number') or number) else ""
    title   = s(finding.get('title'))
    severity = s(finding.get('severity'))

    # Title line with severity badge
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(10)
    nr = p.add_run(f"{num_str}{title}")
    nr.bold = True; nr.font.size = Pt(11); nr.font.color.rgb = DARK; nr.font.name = 'Calibri'
    if severity and severity != '—':
        br = p.add_run(f"  [{severity}]")
        br.font.size = Pt(9); br.font.color.rgb = priority_color(severity)
    if finding.get('module_name'):
        mr = p.add_run(f"  — {finding['module_name']}")
        mr.font.size = Pt(9); mr.font.color.rgb = MUTED

    for label, key in [
        ("Observation",             "observation"),
        ("Why this matters",        "why_it_matters"),
        ("Likely cause",            "likely_cause"),
        ("Business impact",         "impact"),
    ]:
        val = s(finding.get(key), '')
        if val and val != '—':
            labelled(doc, label, val)

    if finding.get('phase'):
        labelled(doc, "Roadmap phase", finding['phase'], value_color=TEAL)


def recommendation_block(doc, rec, number=None):
    num_str = f"Recommendation {rec.get('number') or number}: " if (rec.get('number') or number) else ""
    title   = s(rec.get('title'))

    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(10)
    nr = p.add_run(f"{num_str}{title}")
    nr.bold = True; nr.font.size = Pt(11); nr.font.color.rgb = DARK; nr.font.name = 'Calibri'
    if rec.get('priority'):
        pr = p.add_run(f"  [{s(rec['priority'])} Priority]")
        pr.font.size = Pt(9); pr.font.color.rgb = priority_color(rec['priority'])

    for label, key in [
        ("Recommendation detail",   "detail"),
        ("Rationale",               "rationale"),
        ("Expected benefit",        "expected_benefit"),
    ]:
        val = s(rec.get(key), '')
        if val and val != '—':
            labelled(doc, label, val)


def action_block(doc, action, number=None):
    num_str = f"Action {action.get('number') or number}: " if (action.get('number') or number) else ""
    title   = s(action.get('title'))

    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(10)
    nr = p.add_run(f"{num_str}{title}")
    nr.bold = True; nr.font.size = Pt(11); nr.font.color.rgb = DARK; nr.font.name = 'Calibri'
    if action.get('phase'):
        pr = p.add_run(f"  [{s(action['phase'])}]")
        pr.font.size = Pt(9); pr.font.color.rgb = TEAL

    labelled(doc, "Objective", s(action.get('objective'), ''))

    steps = action.get('key_steps', [])
    if steps:
        sp = doc.add_paragraph()
        sr = sp.add_run("Key steps:")
        sr.bold = True; sr.font.size = Pt(10); sr.font.color.rgb = DARK; sr.font.name = 'Calibri'
        for step in steps:
            bullet(doc, str(step), indent=True)

    labelled(doc, "Suggested owner",    s(action.get('owner_role'), 'Functional Owner'))
    labelled(doc, "Indicative timeline", s(action.get('indicative_timeline'), 'To be confirmed'))
    labelled(doc, "Success indicator",  s(action.get('success_indicator'), ''), value_color=TEAL)


def roadmap_table(doc, items):
    if not items: return
    tbl = doc.add_table(rows=1, cols=5)
    tbl.style = 'Table Grid'
    for cell, lbl in zip(tbl.rows[0].cells, ['Phase', 'Initiative', 'Module', 'Owner', 'Timeline']):
        set_cell_bg(cell, '0F1923')
        r = cell.paragraphs[0].add_run(lbl)
        r.bold = True; r.font.size = Pt(8.5); r.font.color.rgb = WHITE; r.font.name = 'Calibri'

    for item in items:
        row = tbl.add_row()
        phase = s(item.get('phase') or item.get('phase_code'), '—')
        init  = s(item.get('initiative') or item.get('action') or item.get('milestone_name'), '—')
        mod   = s(item.get('module_name'), '—')
        owner = s(item.get('owner_role'), '—')
        tl    = s(item.get('timeline') or item.get('indicative_timeline'), '—')
        ph_code = phase[:2] if len(phase) >= 2 else phase
        ph_r = row.cells[0].paragraphs[0].add_run(phase)
        ph_r.font.size = Pt(8.5)
        ph_r.font.color.rgb = TEAL if ph_code == 'P1' else WARN if ph_code == 'P2' else MUTED
        row.cells[1].paragraphs[0].add_run(init).font.size = Pt(9)
        row.cells[2].paragraphs[0].add_run(mod).font.size = Pt(8.5)
        row.cells[3].paragraphs[0].add_run(owner).font.size = Pt(8.5)
        row.cells[4].paragraphs[0].add_run(tl).font.size = Pt(8.5)
    doc.add_paragraph().paragraph_format.space_after = Pt(6)


# ── DOCX generator ────────────────────────────────────────────────────────────
def generate_docx(payload, output_path):
    ensure_dir(output_path)
    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin    = Cm(2.2)
        section.bottom_margin = Cm(2.2)
        section.left_margin   = Cm(2.8)
        section.right_margin  = Cm(2.8)

    ctx = payload.get('context', {})
    kind = payload.get('report_kind', 'integrated')
    is_integrated = kind == 'integrated'

    # ── Cover ────────────────────────────────────────────────────────────────
    kinto = doc.add_paragraph()
    kinto.alignment = WD_ALIGN_PARAGRAPH.LEFT
    kr = kinto.add_run('KINTO GLOBAL')
    kr.bold = True; kr.font.size = Pt(10); kr.font.color.rgb = TEAL; kr.font.name = 'Calibri'

    doc.add_paragraph()
    tp = doc.add_paragraph()
    tr = tp.add_run(s(payload.get('title'), 'Executive Diagnostic Report'))
    tr.bold = True; tr.font.size = Pt(26); tr.font.color.rgb = DARK; tr.font.name = 'Calibri'
    tp.paragraph_format.space_after = Pt(4)

    sp = doc.add_paragraph()
    sr2 = sp.add_run(s(payload.get('subtitle'), ''))
    sr2.font.size = Pt(11.5); sr2.font.color.rgb = MUTED; sr2.font.name = 'Calibri'

    doc.add_paragraph()
    add_divider(doc, '1ABCB0')

    conf = doc.add_paragraph()
    cr2 = conf.add_run('CONFIDENTIAL  ')
    cr2.bold = True; cr2.font.size = Pt(8); cr2.font.color.rgb = DARK
    cr3 = conf.add_run('This report is prepared solely for the named client and authorised senior stakeholders. Distribution should be controlled accordingly.')
    cr3.font.size = Pt(8); cr3.font.color.rgb = MUTED

    doc.add_page_break()

    # ── 01 Executive Summary ─────────────────────────────────────────────────
    es = payload.get('executive_summary', {})
    section_tag(doc, '01 — Executive Summary')
    h1(doc, 'Executive Summary')

    if es.get('headline'):
        hl = doc.add_paragraph()
        hr = hl.add_run(es['headline'])
        hr.bold = True; hr.font.size = Pt(12); hr.font.color.rgb = DARK
        hl.paragraph_format.space_after = Pt(10)

    if es.get('assessment_context'):
        body(doc, es['assessment_context'])

    if es.get('overall_view'):
        h3(doc, 'Overall Assessment View')
        body(doc, es['overall_view'])

    # Strengths / Weaknesses two-column style
    strengths  = es.get('key_strengths', [])
    weaknesses = es.get('key_weaknesses', [])

    if strengths or weaknesses:
        if strengths:
            h3(doc, 'Key Strengths', color=SUCCESS)
            for item in strengths[:4]: bullet(doc, item, color=SUCCESS)

        if weaknesses:
            h3(doc, 'Key Weaknesses / Priority Gaps', color=DANGER)
            for item in weaknesses[:5]: bullet(doc, item, color=DANGER)

    if es.get('biggest_risks'):
        h3(doc, 'Biggest Risks')
        for item in es['biggest_risks']: bullet(doc, item)

    if es.get('biggest_opportunities'):
        h3(doc, 'Biggest Opportunities', color=TEAL)
        for item in es['biggest_opportunities']: bullet(doc, item, color=TEAL)

    if es.get('top_strategic_priorities'):
        h3(doc, 'Top Strategic Priorities')
        for i, pri in enumerate(es['top_strategic_priorities'][:5], 1):
            p = doc.add_paragraph()
            nr = p.add_run(f'{i}. ')
            nr.bold = True; nr.font.size = Pt(10); nr.font.color.rgb = TEAL
            pr = p.add_run(str(pri))
            pr.font.size = Pt(10); pr.font.color.rgb = DARK

    if es.get('financial_impact'):
        h3(doc, 'Financial Impact')
        body(doc, es['financial_impact'])

    if es.get('conclusion'):
        h3(doc, 'Advisory Conclusion')
        body(doc, es['conclusion'])

    add_divider(doc)
    doc.add_page_break()

    # ── 02 Assessment Scope & Context ────────────────────────────────────────
    section_tag(doc, '02 — Assessment Scope')
    h1(doc, 'Assessment Context and Scope')
    body(doc, f'Client: {s(ctx.get("client_name"))}')
    body(doc, f'Assessment: {s(ctx.get("assessment_name"))}')
    if ctx.get('industry'):  body(doc, f'Industry: {ctx["industry"]}')
    if ctx.get('company_size'): body(doc, f'Company size: {ctx["company_size"]}')
    body(doc, 'This assessment was conducted using the Kinto Global diagnostic methodology, covering five modules: Operational Audit, Revenue Leakage, Data Foundation, AI Readiness, and AI Use Case Prioritisation. Each module applies a structured evaluation framework to surface findings, quantify risk, and define a prioritised improvement roadmap.')
    add_divider(doc)

    # ── 03 Portfolio Scores ───────────────────────────────────────────────────
    module_scores = payload.get('module_scores', [])
    if not module_scores:
        bpo = payload.get('business_performance_overview', {})
        module_scores = bpo.get('scores', [])

    if module_scores:
        section_tag(doc, '03 — Overall Performance')
        h1(doc, 'Business Performance Overview')
        body(doc, 'The table below summarises the maturity score, band, and priority level for each module or domain assessed. Scores reflect the quality of what has been assessed to date. Modules rated Critical or High priority require the most immediate attention and are addressed first in the action plan.')
        score_table(doc, module_scores)
        add_divider(doc)

    # ── 04 Cross-Cutting Themes ───────────────────────────────────────────────
    themes = payload.get('cross_cutting_themes', [])
    if themes:
        section_tag(doc, '04 — Cross-Cutting Themes')
        h1(doc, 'Cross-Cutting Themes')
        body(doc, 'Several themes emerged consistently across multiple modules. These are presented before the module-level analysis because they represent systemic patterns that single-module interventions are unlikely to resolve in isolation.')
        for i, theme in enumerate(themes, 1):
            h2(doc, f"{i}. {s(theme.get('title'))}")
            body(doc, s(theme.get('observation')))
            if theme.get('affected_modules'):
                labelled(doc, 'Affected modules', ', '.join(theme['affected_modules']))
            if theme.get('recommendation'):
                labelled(doc, 'Cross-module recommendation', theme['recommendation'], value_color=TEAL)
        add_divider(doc)
        doc.add_page_break()

    # ── 05 Module-by-Module Analysis ──────────────────────────────────────────
    module_sections = payload.get('module_sections', [])
    if not module_sections and payload.get('module_section'):
        module_sections = [payload['module_section']]

    if module_sections:
        section_tag(doc, '05 — Module Analysis')
        h1(doc, 'Module-by-Module Analysis')
        body(doc, 'The following sections present a detailed analysis for each diagnostic module. Each section covers current state, key findings, recommendations, action plan, and expected outcomes.')

        for ms in module_sections:
            doc.add_page_break()
            mod_name = s(ms.get('module_name'))
            score_val = float(ms.get('score') or 0)
            score_bnd = s(ms.get('score_band'))

            section_tag(doc, f'Module — {mod_name}')
            h1(doc, mod_name, color=DARK)

            # Score indicator
            score_p = doc.add_paragraph()
            sc_r = score_p.add_run(f'Score: {score_val:.1f}%  ')
            sc_r.bold = True; sc_r.font.size = Pt(13); sc_r.font.color.rgb = score_color(score_val)
            bnd_r = score_p.add_run(f'({score_bnd})')
            bnd_r.font.size = Pt(11); bnd_r.font.color.rgb = MUTED

            # Section overview
            h2(doc, '5.1  Section Overview')
            body(doc, ms.get('current_state', ''))

            # Area scores
            if ms.get('area_scores'):
                h2(doc, '5.2  Domain / Area Scores')
                score_table(doc, ms['area_scores'])

            # Current state
            if ms.get('cross_cutting_themes'):
                h2(doc, '5.3  Key Themes Within This Module')
                for theme in ms['cross_cutting_themes']:
                    bullet(doc, str(theme))

            # Findings
            findings = ms.get('findings', [])
            if findings:
                h2(doc, f'5.4  Key Findings ({len(findings)} identified)')
                body(doc, ms.get('transition_note', ''))
                for f in findings:
                    finding_block(doc, f)

            # Recommendations
            recs = ms.get('recommendations', [])
            if recs:
                h2(doc, f'5.5  Recommendations ({len(recs)} recommendations)')
                for r in recs:
                    recommendation_block(doc, r)

            # Actions
            actions = ms.get('actions', [])
            if actions:
                h2(doc, f'5.6  Action Plan ({len(actions)} actions)')
                for a in actions:
                    action_block(doc, a)

            # Metrics
            metrics = ms.get('metrics', [])
            if metrics:
                h2(doc, '5.7  Performance Metrics')
                tbl = doc.add_table(rows=1, cols=5)
                tbl.style = 'Table Grid'
                for cell, lbl in zip(tbl.rows[0].cells, ['Metric', 'Baseline', 'Current', 'Target', 'RAG']):
                    set_cell_bg(cell, '0F1923')
                    r = cell.paragraphs[0].add_run(lbl)
                    r.bold = True; r.font.size = Pt(8.5); r.font.color.rgb = WHITE
                for m in metrics:
                    row = tbl.add_row()
                    row.cells[0].paragraphs[0].add_run(s(m.get('name'))).font.size = Pt(9)
                    row.cells[1].paragraphs[0].add_run(s(m.get('baseline'))).font.size = Pt(9)
                    row.cells[2].paragraphs[0].add_run(s(m.get('current'))).font.size = Pt(9)
                    row.cells[3].paragraphs[0].add_run(s(m.get('target'))).font.size = Pt(9)
                    rag = s(m.get('rag'), '')
                    rr = row.cells[4].paragraphs[0].add_run(rag)
                    rr.font.size = Pt(9)
                    rr.font.color.rgb = DANGER if 'red' in rag.lower() or 'action' in rag.lower() else WARN if 'amber' in rag.lower() else TEAL
                doc.add_paragraph()

            # Expected outcomes
            if ms.get('expected_outcomes'):
                h2(doc, '5.8  Expected Outcomes')
                for outcome in ms['expected_outcomes']:
                    bullet(doc, outcome, color=TEAL)

        add_divider(doc)

    # ── 06 Consolidated Strategic Priorities ──────────────────────────────────
    doc.add_page_break()
    strategic = payload.get('strategic_priorities', {})
    if not strategic:
        flat = payload.get('strategic_priorities_flat') or payload.get('consolidated', {}).get('actions', [])
        if flat:
            strategic = {'all': flat}

    if strategic:
        section_tag(doc, '06 — Strategic Priorities')
        h1(doc, 'Consolidated Strategic Priorities')
        body(doc, 'The following priorities consolidate the highest-impact actions across all modules into a sequenced implementation plan. Priorities are organised by phase, allowing the business to build foundations before advancing to optimisation and AI enablement.')

        for phase_key, phase_label in [
            ('p1', 'Phase 1 — Stabilise & Protect  (0–30 days)'),
            ('p2', 'Phase 2 — Standardise & Strengthen  (30–90 days)'),
            ('p3', 'Phase 3 — Optimise, Automate & AI-Enable  (90+ days)'),
            ('all', 'Priority Actions'),
        ]:
            items = strategic.get(phase_key, [])
            if items:
                h2(doc, phase_label, color=TEAL if phase_key == 'p1' else WARN if phase_key == 'p2' else MUTED)
                for i, item in enumerate(items[:6], 1):
                    if isinstance(item, dict):
                        p = doc.add_paragraph()
                        p.paragraph_format.space_before = Pt(5)
                        ir = p.add_run(f'{i}. {s(item.get("title") or item.get("initiative"))}')
                        ir.bold = True; ir.font.size = Pt(10.5); ir.font.color.rgb = DARK
                        if item.get('objective') or item.get('impact'):
                            bp = doc.add_paragraph()
                            bp.paragraph_format.left_indent = Inches(0.25)
                            br = bp.add_run(s(item.get('objective') or item.get('impact'), ''))
                            br.font.size = Pt(9.5); br.font.color.rgb = MUTED
                        labelled(doc, 'Owner', s(item.get('owner_role') or item.get('priority_score'), ''))
                        labelled(doc, 'Timeline', s(item.get('indicative_timeline') or item.get('effort'), ''))
                    else:
                        bullet(doc, str(item))
        add_divider(doc)

    # ── 06b Ordered Action Plan ───────────────────────────────────────────────
    ordered_plan = payload.get('ordered_plan', {})
    engagement_framing = payload.get('engagement_framing', '')
    if ordered_plan and (ordered_plan.get('immediate') or ordered_plan.get('near_term')):
        section_tag(doc, '06b — Priority Action Plan')
        h1(doc, 'Priority Action Plan')
        if engagement_framing:
            body(doc, engagement_framing)
        body(doc, s(ordered_plan.get('summary'), ''))

        immediate = ordered_plan.get('immediate', [])
        near_term = ordered_plan.get('near_term', [])
        later     = ordered_plan.get('later', [])

        if immediate:
            h2(doc, 'Do First — Phase 1: Stabilise & Protect  (0–30 days)', color=DANGER)
            body(doc, 'These actions address the highest-priority findings. They should be assigned to named owners immediately and tracked weekly.')
            for item in immediate:
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(6)
                nr = p.add_run(f"{item.get('number', '')}. {s(item.get('title'))}")
                nr.bold = True; nr.font.size = Pt(11); nr.font.color.rgb = DARK; nr.font.name = 'Calibri'
                labelled(doc, 'Why', s(item.get('why'), ''))
                labelled(doc, 'Owner', s(item.get('owner'), 'Functional Owner'))
                labelled(doc, 'Timeline', s(item.get('timeline'), '30 days'), value_color=DANGER)

        if near_term:
            h2(doc, 'Do Next — Phase 2: Standardise & Strengthen  (30–90 days)', color=WARN)
            body(doc, 'These actions build on the Phase 1 improvements and address the next tier of findings.')
            for item in near_term:
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(6)
                nr = p.add_run(f"{item.get('number', '')}. {s(item.get('title'))}")
                nr.bold = True; nr.font.size = Pt(11); nr.font.color.rgb = DARK; nr.font.name = 'Calibri'
                labelled(doc, 'Why', s(item.get('why'), ''))
                labelled(doc, 'Owner', s(item.get('owner'), 'Functional Owner'))
                labelled(doc, 'Timeline', s(item.get('timeline'), '30–90 days'), value_color=WARN)

        if later:
            h2(doc, 'Phase 3: Optimise, Automate & AI-Enable  (90+ days)', color=MUTED)
            for item in later:
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(4)
                nr = p.add_run(f"{item.get('number', '')}. {s(item.get('title'))}")
                nr.bold = True; nr.font.size = Pt(10.5); nr.font.color.rgb = MUTED; nr.font.name = 'Calibri'
                labelled(doc, 'Owner', s(item.get('owner'), ''))
        add_divider(doc)

    # ── 07 Implementation Roadmap ─────────────────────────────────────────────
    doc.add_page_break()
    roadmap = payload.get('consolidated', {}).get('roadmap', []) or payload.get('implementation_roadmap', [])
    if roadmap:
        section_tag(doc, '07 — Implementation Roadmap')
        h1(doc, 'Implementation Roadmap')
        body(doc, 'The roadmap below presents all priority initiatives organised by implementation phase. Each initiative is linked to a specific module finding, carries a recommended owner role, and includes an indicative timeline. Progress against this roadmap should be reviewed on a monthly basis.')
        roadmap_table(doc, roadmap[:20])
        add_divider(doc)

    # ── 08 Expected Outcomes ──────────────────────────────────────────────────
    outcomes = payload.get('expected_outcomes', [])
    if outcomes:
        section_tag(doc, '08 — Expected Outcomes')
        h1(doc, 'Expected Business Outcomes')
        body(doc, 'If the priority actions are executed with appropriate ownership and governance, the following outcomes should be realised within the indicated timeframes.')
        for outcome in outcomes:
            bullet(doc, str(outcome), color=TEAL)
        add_divider(doc)

    # ── 09 Closing Advisory Note ──────────────────────────────────────────────
    closing = payload.get('closing_note', {})
    if closing:
        doc.add_page_break()
        section_tag(doc, '09 — Closing Note')
        h1(doc, s(closing.get('title'), 'Closing Advisory Note'))
        body(doc, closing.get('message', ''), size=11)

    # ── Footer ────────────────────────────────────────────────────────────────
    doc.add_paragraph()
    add_divider(doc)
    fp = doc.add_paragraph()
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fr = fp.add_run('KINTO GLOBAL  ·  Visionary Minds, Ingenious Designs  ·  Confidential')
    fr.font.size = Pt(8); fr.font.color.rgb = MUTED

    doc.save(output_path)


# ── PPTX generator ────────────────────────────────────────────────────────────
def pptx_header(slide, title, tag, W, H):
    bar = slide.shapes.add_shape(1, 0, 0, W, Pi(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = P_DARK; bar.line.fill.background()
    tf = bar.text_frame; tf.margin_top = Pi(0.15); tf.margin_left = Pi(0.35)
    p = tf.paragraphs[0]
    if tag:
        r1 = p.add_run(); r1.text = f'{tag}  '; r1.font.size = PPt(9); r1.font.color.rgb = P_TEAL; r1.font.bold = True
    r2 = p.add_run(); r2.text = str(title); r2.font.size = PPt(16); r2.font.color.rgb = P_WHITE; r2.font.bold = True


def pptx_body(slide, text, x, y, w, h, size=12, color=None, bold=False, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = wrap; tf.auto_size = None
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = str(text)[:600]
    r.font.size = PPt(size); r.font.bold = bold
    r.font.color.rgb = color or P_MUTED


def pptx_bullets(slide, items, x, y, w, h, size=11, color=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = f'• {str(item)[:200]}'
        r.font.size = PPt(size)
        r.font.color.rgb = color or P_MUTED


def generate_pptx(payload, output_path):
    ensure_dir(output_path)
    prs = Presentation()
    prs.slide_width  = Pi(13.33)
    prs.slide_height = Pi(7.5)
    W, H = prs.slide_width, prs.slide_height
    SL = prs.slide_layouts[6]  # blank

    def new_slide():
        return prs.slides.add_slide(SL)

    ctx = payload.get('context', {})
    es  = payload.get('executive_summary', {})

    # ── Slide 1: Cover ────────────────────────────────────────────────────────
    slide = new_slide()
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = P_DARK; bg.line.fill.background()
    accent = slide.shapes.add_shape(1, 0, Pi(6.1), W, Pi(0.06))
    accent.fill.solid(); accent.fill.fore_color.rgb = P_TEAL; accent.line.fill.background()
    pptx_body(slide, 'KINTO GLOBAL', Pi(0.55), Pi(0.35), Pi(5), Pi(0.5), 11, P_TEAL, bold=True)
    pptx_body(slide, s(payload.get('title')), Pi(0.55), Pi(1.8), Pi(12.2), Pi(2.2), 30, P_WHITE, bold=True)
    pptx_body(slide, s(payload.get('subtitle')), Pi(0.55), Pi(4.2), Pi(10), Pi(0.7), 13, P_MUTED)
    pptx_body(slide, 'CONFIDENTIAL', Pi(0.55), Pi(6.9), Pi(6), Pi(0.4), 9, P_MUTED, bold=True)

    # ── Slide 2: Executive Summary ────────────────────────────────────────────
    slide = new_slide()
    pptx_header(slide, 'Executive Summary', '01', W, H)
    y = Pi(1.1)
    headline = s(es.get('headline'), '')
    if headline:
        pptx_body(slide, headline, Pi(0.4), y, Pi(12.5), Pi(1.6), 13, P_DARK, bold=True)
        y += Pi(1.7)
    context_txt = s(es.get('assessment_context') or es.get('overall_view'), '')
    if context_txt:
        pptx_body(slide, context_txt[:400], Pi(0.4), y, Pi(12.5), Pi(1.8), 10.5, P_MUTED)
        y += Pi(1.9)
    if es.get('financial_impact'):
        pptx_body(slide, f"Financial Impact:  {s(es['financial_impact'])}", Pi(0.4), y, Pi(12.5), Pi(0.7), 11, P_DANGER, bold=True)

    # ── Slide 3: Key Risks & Opportunities ───────────────────────────────────
    slide = new_slide()
    pptx_header(slide, 'Key Risks & Opportunities', '02', W, H)
    risks = es.get('biggest_risks', [])
    opps  = es.get('biggest_opportunities', [])
    pptx_body(slide, 'BIGGEST RISKS', Pi(0.4), Pi(1.1), Pi(6), Pi(0.4), 9, P_DANGER, bold=True)
    pptx_bullets(slide, risks[:4], Pi(0.4), Pi(1.55), Pi(6), Pi(2.8), 11, P_DARK)
    pptx_body(slide, 'BIGGEST OPPORTUNITIES', Pi(6.8), Pi(1.1), Pi(6), Pi(0.4), 9, P_TEAL, bold=True)
    pptx_bullets(slide, opps[:4], Pi(6.8), Pi(1.55), Pi(6), Pi(2.8), 11, P_DARK)
    if es.get('conclusion'):
        pptx_body(slide, s(es['conclusion'])[:300], Pi(0.4), Pi(5.5), Pi(12.5), Pi(1.5), 10, P_MUTED)

    # ── Slide 4: Business Performance Scores ─────────────────────────────────
    scores = payload.get('module_scores', []) or payload.get('business_performance_overview', {}).get('scores', [])
    if scores:
        slide = new_slide()
        pptx_header(slide, 'Business Performance Overview', '03', W, H)
        col_count = min(len(scores[:5]), 5)
        col_w = Pi(12.5) / col_count
        for i, sc in enumerate(scores[:5]):
            pct_val = float(sc.get('score') or sc.get('percentage') or 0)
            name    = s(sc.get('module_name') or sc.get('name'))
            bnd     = s(sc.get('score_band') or sc.get('band'))
            pri     = s(sc.get('priority', ''))
            col_color = p_score_color(pct_val)
            x = Pi(0.4) + i * col_w
            # Card bg
            card = slide.shapes.add_shape(1, x, Pi(1.1), col_w - Pi(0.1), Pi(5.8))
            card.fill.solid(); card.fill.fore_color.rgb = PRGBColor(248, 250, 252)
            card.line.color.rgb = PRGBColor(226, 232, 240)
            # Module name
            nb = slide.shapes.add_textbox(x + Pi(0.1), Pi(1.2), col_w - Pi(0.2), Pi(0.9))
            tf = nb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = name; r.font.size = PPt(10); r.font.color.rgb = P_DARK; r.font.bold = True
            # Score
            pb = slide.shapes.add_textbox(x + Pi(0.05), Pi(2.2), col_w - Pi(0.1), Pi(1.0))
            tf = pb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = f'{pct_val:.0f}%'; r.font.size = PPt(30); r.font.color.rgb = col_color; r.font.bold = True
            # Band
            bb = slide.shapes.add_textbox(x + Pi(0.05), Pi(3.3), col_w - Pi(0.1), Pi(0.5))
            tf = bb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = bnd; r.font.size = PPt(9.5); r.font.color.rgb = col_color
            # Priority
            if pri:
                prb = slide.shapes.add_textbox(x + Pi(0.05), Pi(3.9), col_w - Pi(0.1), Pi(0.4))
                tf = prb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = f'{pri} Priority'; r.font.size = PPt(8.5); r.font.color.rgb = p_priority_color(pri)

    # ── Slide 5+: Module Summaries ────────────────────────────────────────────
    module_sections = payload.get('module_sections', [])
    if not module_sections and payload.get('module_section'):
        module_sections = [payload['module_section']]

    for ms in module_sections[:5]:
        slide = new_slide()
        mod_name  = s(ms.get('module_name'))
        score_val = float(ms.get('score') or 0)
        col       = p_score_color(score_val)
        pptx_header(slide, f'{mod_name} — {score_val:.0f}% ({s(ms.get("score_band"))})', '05', W, H)

        # Current state
        current = s(ms.get('current_state', ''))[:400]
        if current:
            pptx_body(slide, current, Pi(0.4), Pi(1.1), Pi(8), Pi(2.2), 10.5, P_MUTED)

        # Top 3 findings
        findings = ms.get('findings', [])[:3]
        if findings:
            fb = slide.shapes.add_textbox(Pi(8.6), Pi(1.1), Pi(4.5), Pi(2.8))
            tf = fb.text_frame; tf.word_wrap = True
            htf = tf.paragraphs[0]
            hr = htf.add_run(); hr.text = 'Key Findings:'; hr.font.size = PPt(9); hr.font.bold = True; hr.font.color.rgb = P_DARK
            for f in findings:
                p2 = tf.add_paragraph()
                r = p2.add_run(); r.text = f'• {s(f.get("title"))[:80]}'; r.font.size = PPt(9.5); r.font.color.rgb = P_DARK

        # Top 3 actions
        actions = ms.get('actions', [])[:3]
        if actions:
            ab = slide.shapes.add_textbox(Pi(0.4), Pi(3.5), Pi(12.5), Pi(2.5))
            tf = ab.text_frame; tf.word_wrap = True
            htf = tf.paragraphs[0]
            hr = htf.add_run(); hr.text = 'Priority Actions:'; hr.font.size = PPt(9); hr.font.bold = True; hr.font.color.rgb = P_TEAL
            for a in actions:
                p2 = tf.add_paragraph()
                r = p2.add_run(); r.text = f'▸  {s(a.get("title"))} — {s(a.get("owner_role",""))} · {s(a.get("indicative_timeline",""))}'; r.font.size = PPt(9.5); r.font.color.rgb = P_DARK

    # ── Slide: Priority Action Plan ───────────────────────────────────────────
    ordered_plan = payload.get('ordered_plan', {})
    if ordered_plan and (ordered_plan.get('immediate') or ordered_plan.get('near_term')):
        immediate = ordered_plan.get('immediate', [])
        near_term = ordered_plan.get('near_term', [])
        slide = new_slide()
        pptx_header(slide, 'Priority Action Plan', '06', W, H)

        plan_summary = ordered_plan.get('summary', '')
        if plan_summary:
            pptx_body(slide, plan_summary[:250], Pi(0.4), Pi(1.1), Pi(12.5), Pi(0.8), 10.5, P_MUTED)

        col_gap = Pi(6.4)
        # P1 column
        if immediate:
            pptx_body(slide, 'DO FIRST — Phase 1 (0–30 days)', Pi(0.4), Pi(1.95), Pi(6), Pi(0.4), 9, P_DANGER, bold=True)
            box = slide.shapes.add_textbox(Pi(0.4), Pi(2.4), Pi(6), Pi(4.6))
            tf = box.text_frame; tf.word_wrap = True
            for i, item in enumerate(immediate[:5]):
                p2 = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                r = p2.add_run()
                r.text = f"{item.get('number', i+1)}.  {s(item.get('title'))[:80]}"
                r.font.size = PPt(10); r.font.color.rgb = P_DARK; r.font.bold = True
                owner = s(item.get('owner'), '')
                if owner and owner != '—':
                    p3 = tf.add_paragraph()
                    r3 = p3.add_run(); r3.text = f"    Owner: {owner}"; r3.font.size = PPt(8.5); r3.font.color.rgb = P_MUTED
        # P2 column
        if near_term:
            pptx_body(slide, 'DO NEXT — Phase 2 (30–90 days)', Pi(0.4) + col_gap, Pi(1.95), Pi(6), Pi(0.4), 9, P_WARN, bold=True)
            box2 = slide.shapes.add_textbox(Pi(0.4) + col_gap, Pi(2.4), Pi(6), Pi(4.6))
            tf2 = box2.text_frame; tf2.word_wrap = True
            for i, item in enumerate(near_term[:5]):
                p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                r = p2.add_run()
                r.text = f"{item.get('number', len(immediate)+i+1)}.  {s(item.get('title'))[:80]}"
                r.font.size = PPt(10); r.font.color.rgb = P_DARK; r.font.bold = True

    # ── Slide: Strategic Priorities ───────────────────────────────────────────
    strategic = payload.get('strategic_priorities', {})
    flat_actions = payload.get('consolidated', {}).get('actions', []) or payload.get('strategic_priorities_flat', [])

    slide = new_slide()
    pptx_header(slide, 'Strategic Priorities by Phase', '06', W, H)
    for i, (phase_key, phase_label, col) in enumerate([
        ('p1', 'Phase 1 — Stabilise & Protect', P_DANGER),
        ('p2', 'Phase 2 — Standardise & Strengthen', P_WARN),
        ('p3', 'Phase 3 — Optimise & AI-Enable', P_TEAL),
    ]):
        items = strategic.get(phase_key, []) or flat_actions[i*3:(i+1)*3]
        if not items: continue
        x = Pi(0.4) + i * Pi(4.3)
        pptx_body(slide, phase_label, x, Pi(1.1), Pi(4.2), Pi(0.5), 9, col, bold=True)
        bx = slide.shapes.add_textbox(x, Pi(1.65), Pi(4.1), Pi(5.3))
        tf = bx.text_frame; tf.word_wrap = True
        for j, item in enumerate(items[:5]):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            title = s(item.get('title') or item.get('initiative'), '') if isinstance(item, dict) else str(item)
            r = p.add_run(); r.text = f'{j+1}. {title[:100]}'; r.font.size = PPt(10); r.font.color.rgb = P_DARK

    # ── Slide: Implementation Roadmap ─────────────────────────────────────────
    roadmap = payload.get('consolidated', {}).get('roadmap', []) or payload.get('implementation_roadmap', [])
    if roadmap:
        slide = new_slide()
        pptx_header(slide, 'Implementation Roadmap', '07', W, H)
        box = slide.shapes.add_textbox(Pi(0.4), Pi(1.1), Pi(12.5), Pi(5.8))
        tf = box.text_frame; tf.word_wrap = True
        phases = {}
        for r in roadmap[:18]:
            ph = s(r.get('phase') or r.get('phase_code'), 'Unphased')
            phases.setdefault(ph, []).append(r)
        first = True
        for phase, items in list(phases.items())[:3]:
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            col = P_TEAL if 'P1' in phase or 'Stabilise' in phase else P_WARN if 'P2' in phase or 'Standard' in phase else P_MUTED
            r = p.add_run(); r.text = phase; r.font.size = PPt(11); r.font.color.rgb = col; r.font.bold = True
            for item in items[:4]:
                p2 = tf.add_paragraph()
                action = s(item.get('initiative') or item.get('action') or item.get('milestone_name'))
                owner  = s(item.get('owner_role'), '')
                r2 = p2.add_run(); r2.text = f'   • {action}{(" — " + owner) if owner else ""}'; r2.font.size = PPt(10); r2.font.color.rgb = P_DARK

    # ── Slide: AI Opportunities ────────────────────────────────────────────────
    ai_opps = payload.get('ai_opportunities', [])
    if ai_opps:
        slide = new_slide()
        pptx_header(slide, 'AI & Automation Opportunities', '08', W, H)
        box = slide.shapes.add_textbox(Pi(0.4), Pi(1.1), Pi(12.5), Pi(5.8))
        tf = box.text_frame; tf.word_wrap = True
        for i, opp in enumerate(ai_opps[:6]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            uc = s(opp.get('use_case')); oc = s(opp.get('outcome'), '')
            r1 = p.add_run(); r1.text = f'⚡  {uc}'; r1.font.size = PPt(12); r1.font.color.rgb = P_DARK; r1.font.bold = True
            if oc and oc != '—':
                p2 = tf.add_paragraph()
                r2 = p2.add_run(); r2.text = f'    {oc[:180]}'; r2.font.size = PPt(10); r2.font.color.rgb = P_MUTED

    # ── Slide: Expected Outcomes ───────────────────────────────────────────────
    outcomes = payload.get('expected_outcomes', [])
    if outcomes:
        slide = new_slide()
        pptx_header(slide, 'Expected Business Outcomes', '09', W, H)
        pptx_bullets(slide, [str(o)[:200] for o in outcomes[:7]], Pi(0.4), Pi(1.1), Pi(12.5), Pi(5.5), 12, P_TEAL)

    # ── Closing slide ─────────────────────────────────────────────────────────
    slide = new_slide()
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = P_DARK; bg.line.fill.background()
    accent = slide.shapes.add_shape(1, 0, Pi(3.7), W, Pi(0.06))
    accent.fill.solid(); accent.fill.fore_color.rgb = P_TEAL; accent.line.fill.background()
    closing_box = slide.shapes.add_textbox(Pi(0.5), Pi(2.7), Pi(12.3), Pi(1.0))
    tf = closing_box.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'KINTO GLOBAL'; r.font.size = PPt(20); r.font.color.rgb = P_WHITE; r.font.bold = True
    tag_box2 = slide.shapes.add_textbox(Pi(0.5), Pi(3.85), Pi(12.3), Pi(0.6))
    tf = tag_box2.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'Visionary Minds, Ingenious Designs'; r.font.size = PPt(12); r.font.color.rgb = P_TEAL

    prs.save(output_path)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--payload', required=True)
    parser.add_argument('--docx',    required=True)
    parser.add_argument('--pptx',    required=True)
    args = parser.parse_args()
    with open(args.payload, 'r', encoding='utf-8') as fh:
        payload = json.load(fh)
    generate_docx(payload, args.docx)
    generate_pptx(payload, args.pptx)
    print(f'Generated: {args.docx}')
    print(f'Generated: {args.pptx}')


if __name__ == '__main__':
    main()
